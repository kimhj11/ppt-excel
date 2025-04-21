import streamlit as st
from pptx import Presentation
import pandas as pd
import re
import io
import os

st.set_page_config(page_title="PPT to Excel 변환기", layout="wide")
st.title("✍️ 온라인 평가문항 양식 자동 변환하기")

uploaded_files = st.file_uploader("📤 PPT 파일 업로드 (.pptx만 가능)", type=["pptx"], accept_multiple_files=True)
base_excel = st.file_uploader("📂 기존 결과 엑셀 업로드 (선택)", type=["xlsx"])
set_number = st.number_input("📦 세트 번호", min_value=1, value=1, step=1)

data_rows = []
failed_slides = {}

if uploaded_files:
    st.info("각 파일마다 추출할 슬라이드 번호를 입력하세요 (예: 3,5,7 또는 전체).")
    slide_inputs = {}
    for file in uploaded_files:
        slide_input = st.text_input(f"{file.name} 슬라이드 번호", key=f"slide_{file.name}")
        slide_inputs[file.name] = slide_input

    with st.spinner("🔍 문항을 추출 중입니다..."):
        for file in uploaded_files:
            filename = file.name
            prs = Presentation(file)

            slide_input = slide_inputs.get(filename, "")
            if slide_input.strip().lower() == "전체":
                target_slides = range(len(prs.slides))
            else:
                try:
                    target_slides = [int(n.strip()) - 1 for n in slide_input.split(",")]
                except:
                    st.error(f"{filename}의 슬라이드 번호 형식이 잘못되었습니다.")
                    continue

            lesson_match = re.search(r"(\d+)차시", filename)
            lesson_name = lesson_match.group(1) if lesson_match else "1"

            for idx in target_slides:
                if idx >= len(prs.slides): continue
                slide = prs.slides[idx]
                texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            row_text = " ".join(cell.text.strip() for cell in row.cells)
                            texts.append(row_text)

                slide_text = "\n".join(texts)

                with st.expander(f"📝 슬라이드 {idx + 1} ({filename}) 보기 및 수정", expanded=False):
                    user_text = st.text_area(f"✏️ 수정 영역", value=slide_text, height=300, key=f"{filename}_{idx}")

                    pattern = re.compile(
                        r"(?P<번호>\d+)\.\s*(?P<문제>.*?)"
                        r"(?:\n(?P<보기>(?:[①②③④].*?\n?)+))?"
                        r"\n정답\s*[:：]?\s*(?P<정답>[OX①②③④])"
                        r"(?:\n난이도\s*[:：]?\s*(?P<난이도>.*?))?"
                        r"\n해[설석]?\s*[:：]?\s*(?P<해설>.*?)(?=\n\d+\.|\s*\Z)",
                        re.DOTALL
                    )

                    match_count = 0

                    for match in pattern.finditer(user_text):
                        match_count += 1
                        문제 = match.group("문제").strip()
                        정답 = match.group("정답").strip()
                        난이도 = match.group("난이도").strip() if match.group("난이도") else ""
                        해설 = match.group("해설").strip()
                        문항유형 = "OX형" if 정답 in ["O", "X"] else "객관식단일형"

                        정답변환 = {"①": "1", "②": "2", "③": "3", "④": "4"}
                        정답 = 정답변환.get(정답, 정답)

                        보기1 = 보기2 = 보기3 = 보기4 = ""
                        보기_raw = match.group("보기")
                        if 보기_raw and 문항유형 == "객관식단일형":
                            try:
                                보기_split = re.split(r"[①②③④]\s*", 보기_raw)
                                보기_split = [v.strip() for v in 보기_split if v.strip()]
                            except:
                                보기_split = []

                            if len(보기_split) >= 1: 보기1 = 보기_split[0]
                            if len(보기_split) >= 2: 보기2 = 보기_split[1]
                            if len(보기_split) >= 3: 보기3 = 보기_split[2]
                            if len(보기_split) >= 4: 보기4 = 보기_split[3]

                        data_rows.append({
                            "문항유형": 문항유형,
                            "종류": "텍스트",
                            "난이도": 난이도,
                            "문제": 문제,
                            "정답": 정답,
                            "보기①": 보기1,
                            "보기②": 보기2,
                            "보기③": 보기3,
                            "보기④": 보기4,
                            "해설": 해설,
                            "세트": set_number,
                            "차시": lesson_name
                        })

                    if match_count == 0:
                        failed_slides.setdefault(filename, []).append(idx + 1)

if data_rows:
    df_new = pd.DataFrame(data_rows)
    df_new.insert(0, "번호", range(1, len(df_new) + 1))

    if base_excel:
        df_old = pd.read_excel(base_excel)
        df_combined = pd.concat([df_old, df_new], ignore_index=True)
    else:
        df_combined = df_new

    st.subheader("📄 미리보기")
    st.dataframe(df_combined)

    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df_combined.to_excel(writer, index=False, sheet_name="문항")

    st.download_button(
        label="📥 엑셀 다운로드",
        data=output.getvalue(),
        file_name="온라인_평가문항_최종결과.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )

    if failed_slides:
        st.warning("⚠️ 추출에 실패한 슬라이드가 있습니다:")
        for fname, slide_nums in failed_slides.items():
            st.markdown(f"- **{fname}**: 슬라이드 {', '.join(map(str, slide_nums))}")
else:
    st.info("📤 PPT 업로드 후 슬라이드 번호를 입력하면 추출 결과가 여기에 표시됩니다.")
